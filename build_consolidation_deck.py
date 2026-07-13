"""
Build an informational PowerPoint consolidation analysis deck.
Run: python build_consolidation_deck.py
Output: ietf_consolidation_analysis.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color palette ──────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x3A, 0x5C)
BLUE   = RGBColor(0x25, 0x6E, 0xB4)   # WG-adopted / absorption-eligible
TEAL   = RGBColor(0x00, 0x82, 0x80)   # companion / closed window
ORANGE = RGBColor(0xC4, 0x62, 0x00)   # no WG anchor
RED    = RGBColor(0xB3, 0x12, 0x12)   # blocking / expired
LGREY  = RGBColor(0xF4, 0xF6, 0xF9)
DGREY  = RGBColor(0x3C, 0x3C, 0x3C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MIDBLUE= RGBColor(0xE8, 0xF0, 0xFB)
SILVER = RGBColor(0xDD, 0xE3, 0xEA)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helpers ────────────────────────────────────────────────────────────────────

def slide_bg(slide, color=LGREY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def add_label(slide, text, left, top, width, height,
              size=11, bold=False, color=DGREY,
              align=PP_ALIGN.LEFT, italic=False):
    txbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txbox

def title_bar(slide, text, subtitle=None):
    add_rect(slide, 0, 0, 13.333, 1.1, fill=NAVY)
    add_label(slide, text, 0.35, 0.1, 12.5, 0.75, size=24, bold=True, color=WHITE)
    if subtitle:
        add_label(slide, subtitle, 0.35, 0.72, 12.5, 0.4,
                  size=12, color=RGBColor(0xB8, 0xD0, 0xED))

def footer(slide, text="IETF Agent Authorization — Cluster Analysis  |  July 2026"):
    add_rect(slide, 0, 7.15, 13.333, 0.35, fill=NAVY)
    add_label(slide, text, 0.35, 7.17, 12.6, 0.3,
              size=8, color=RGBColor(0xB8, 0xD0, 0xED))

def draft_row(slide, top, name, description, accent=BLUE, expired=False):
    """One draft entry: colored left rule, bold name, description below."""
    add_rect(slide, 0.4, top, 0.06, 0.72, fill=RED if expired else accent)
    color = RGBColor(0x99, 0x99, 0x99) if expired else accent
    label = f"{name}  [EXPIRED]" if expired else name
    add_label(slide, label, 0.58, top, 12.3, 0.3,
              size=10.5, bold=True, color=color)
    add_label(slide, description, 0.58, top + 0.3, 12.3, 0.44,
              size=10, color=RGBColor(0x77, 0x77, 0x77) if expired else DGREY)

# ── Slide 1: Title ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s, NAVY)
add_label(s, "IETF Agent Authorization", 1.0, 1.6, 11.3, 1.0,
          size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_label(s, "Cluster Analysis", 1.0, 2.55, 11.3, 0.8,
          size=28, color=RGBColor(0xB8, 0xD0, 0xED), align=PP_ALIGN.CENTER)
add_label(s, "Mapping the individual submission landscape by problem space",
          1.5, 3.55, 10.3, 0.6,
          size=15, color=RGBColor(0x9A, 0xBC, 0xD8), align=PP_ALIGN.CENTER, italic=True)
add_label(s, "July 2026  ·  99 Active IETF Drafts  ·  OAuth WG + WIMSE WG",
          1.5, 6.6, 10.3, 0.45,
          size=11, color=RGBColor(0x78, 0xA0, 0xBB), align=PP_ALIGN.CENTER)

# ── Slide 2: WG-Adopted Anchors ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_bar(s, "WG-Adopted Anchors",
          "Eight OAuth WG work items serve as the structural anchors for this landscape")
footer(s)

wg_rows = [
    ("draft-ietf-oauth-transaction-tokens",             "WG, -09 · WGLC"),
    ("draft-ietf-oauth-identity-chaining",              "IESG-approved, RFC Editor queue"),
    ("draft-ietf-oauth-identity-assertion-authz-grant", "WG, -03"),
    ("draft-ietf-oauth-attestation-based-client-auth",  "WG, -10"),
    ("draft-ietf-oauth-first-party-apps",               "WG, -04"),
    ("draft-ietf-oauth-spiffe-client-auth",             "WG (adopted from schwenkschuster), -01"),
    ("draft-ietf-oauth-browser-based-apps",             "WG, -26"),
    ("draft-ietf-oauth-security-topics-update",         "WG, -01"),
]

add_rect(s, 0.4, 1.25, 12.5, 0.38, fill=MIDBLUE)
add_label(s, "Draft", 0.5, 1.28, 7.5, 0.32, size=11, bold=True, color=NAVY)
add_label(s, "Status", 8.1, 1.28, 4.7, 0.32, size=11, bold=True, color=NAVY)

for i, (draft, status) in enumerate(wg_rows):
    top = 1.63 + i * 0.53
    bg = WHITE if i % 2 == 0 else RGBColor(0xF9, 0xFB, 0xFD)
    add_rect(s, 0.4, top, 12.5, 0.51, fill=bg, line=SILVER)
    closed = "RFC Editor" in status
    add_label(s, draft, 0.5, top + 0.08, 7.5, 0.35,
              size=10.5, bold=False, color=TEAL if closed else BLUE)
    add_label(s, status, 8.1, top + 0.08, 4.7, 0.35,
              size=10, color=TEAL if closed else DGREY)

add_label(s,
    "WIMSE WG has its own adopted core suite (arch, identifier, mutual-tls, http-signature) — treated as a separate WG throughout.",
    0.4, 6.58, 12.5, 0.38, size=9, italic=True, color=RGBColor(0x88, 0x88, 0x88))

# ── Slide 3: Cluster Overview ──────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_bar(s, "Nine Clusters", "All 99 active drafts grouped by target problem space")
footer(s)

clusters = [
    ("A", "Transaction Token extensions",  "1 (txn-tokens, WGLC)",        BLUE),
    ("B", "Cross-domain identity",         "1 (identity-chaining, closed)",TEAL),
    ("C", "Client / instance auth",        "1 (attestation-based)",        BLUE),
    ("D", "Actor chain representation",    "0",                            ORANGE),
    ("E", "Token attenuation",             "0",                            ORANGE),
    ("F", "Pre-action permit",             "0",                            ORANGE),
    ("G", "HTTP challenge signals",        "0",                            ORANGE),
    ("H", "AIP name collision",            "0  —  BLOCKING",               RED),
    ("I", "WIMSE agent extensions",        "4 (WIMSE core, separate WG)",  BLUE),
]

add_rect(s, 0.3, 1.22, 12.7, 0.38, fill=MIDBLUE)
for txt, lft, wid in [("", 0.38, 0.6), ("Cluster", 1.05, 4.5), ("WG anchors", 5.65, 7.3)]:
    add_label(s, txt, lft, 1.25, wid, 0.32, size=10, bold=True, color=NAVY)
add_label(s, "Cluster", 1.05, 1.25, 4.5, 0.32, size=10, bold=True, color=NAVY)
add_label(s, "WG anchors in this space", 5.65, 1.25, 7.3, 0.32, size=10, bold=True, color=NAVY)

for i, (letter, topic, anchors, col) in enumerate(clusters):
    top = 1.6 + i * 0.54
    bg = WHITE if i % 2 == 0 else RGBColor(0xF9, 0xFB, 0xFD)
    add_rect(s, 0.3, top, 12.7, 0.52, fill=bg, line=SILVER)
    add_rect(s, 0.3, top, 0.6, 0.52, fill=col)
    add_label(s, letter, 0.3, top + 0.1, 0.6, 0.32,
              size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_label(s, topic,   1.05, top + 0.1, 4.5, 0.32, size=10.5, color=DGREY)
    add_label(s, anchors, 5.65, top + 0.1, 7.3, 0.32, size=10,   color=col)

# ── Cluster slide builder ──────────────────────────────────────────────────────

def cluster_slide(letter, title, subtitle, accent, drafts):
    """
    drafts: list of (name, description, expired=False) tuples.
    Lays out rows automatically; accent sets the left-rule color for non-expired.
    """
    s = prs.slides.add_slide(BLANK)
    slide_bg(s)
    title_bar(s, f"Cluster {letter} — {title}", subtitle)
    footer(s)

    row_h = 0.8          # height per draft row
    gap   = 0.1          # gap between rows
    start = 1.25

    for i, entry in enumerate(drafts):
        name = entry[0]
        desc = entry[1]
        expired = entry[2] if len(entry) > 2 else False
        top = start + i * (row_h + gap)
        draft_row(s, top, name, desc, accent=accent, expired=expired)
        # light separator
        if i < len(drafts) - 1:
            add_rect(s, 0.4, top + row_h, 12.5, 0.01, fill=SILVER)

    return s

# ── Slide 4: Cluster A ─────────────────────────────────────────────────────────
cluster_slide("A",
    "Transaction Token Extensions",
    "WG anchor: draft-ietf-oauth-transaction-tokens (WGLC, -09)",
    BLUE,
    [
        ("draft-araut-oauth-transaction-tokens-for-agents",
         "Profiles TxnTokens with three new claims — act (delegation actor), agentic_ctx (task metadata), and actchain (delegation chain) — giving resource servers a richer picture of which agent acted on whose behalf and in what context."),
        ("draft-liu-oauth-a2a-profile",
         "Profiles TxnTokens specifically for agent-to-agent communication: embeds call-chain context to preserve agent identity, authorization, and operational flow across workloads in trusted environments. Same Huawei author pair as draft-ni-oauth-batch-authorization-delegation.",
         True),  # expired
        ("draft-zhu-oauth-async-delegation",
         "Introduces a Delegation Handle — a short-lived, sender-constrained token that lets an agent refresh its authorization asynchronously without re-prompting the user. Addresses the gap where TxnTokens cover synchronous flows but not long-running or background agent tasks."),
        ("draft-fletcher-transaction-token-chaining-profile",
         "Uses a TxnToken as the subject token in an RFC 8693 Token Exchange to produce a JWT Authorization Grant that crosses trust-domain boundaries. Distinct problem from the WG draft: the WG scope is single-domain; this profile handles cross-domain traversal."),
    ]
)

# ── Slide 5: Cluster B ─────────────────────────────────────────────────────────
cluster_slide("B",
    "Cross-Domain Identity / Token Exchange",
    "WG anchor: draft-ietf-oauth-identity-chaining (IESG-approved, RFC Editor queue)",
    TEAL,
    [
        ("draft-mcguinness-oauth-id-assertion-framework",
         "Defines an Authority Delegation Model and an Issuer Trust Policy JSON document that specifies which issuers are authorized to assert identity into a given namespace. Provides the trust framework that identity-chaining assumes but does not specify."),
        ("draft-mcguinness-oauth-domain-authorized-issuer",
         "Implements the id-assertion-framework trust method using DNS TXT records at _oauth-issuer-policy to publish namespace authorization. Allows relying parties to discover which issuers are DNS-authorized for a domain without prior out-of-band configuration."),
        ("draft-mcguinness-token-xchg-target-svc-disco",
         "Defines how a client discovers the correct token-exchange endpoint for a target service. Fills the operational gap where identity-chaining defines the exchange mechanics but not how the client learns where to send the request."),
        ("draft-fletcher-transaction-token-chaining-profile",
         "Appears in both Cluster A and B: it composes identity-chaining with TxnTokens to enable cross-domain agent delegation while preserving the full delegation chain in the resulting token."),
        ("draft-liu-oauth-chain-delegation",
         "Defines a delegation_chain array carrying per-hop authorization constraints within the token itself. Addresses cross-domain delegation lineage alongside identity-chaining's identity propagation."),
    ]
)

# ── Slide 6: Cluster C ─────────────────────────────────────────────────────────
cluster_slide("C",
    "Client Authentication / Instance Identity",
    "WG anchor: draft-ietf-oauth-attestation-based-client-auth (WG, -10)",
    BLUE,
    [
        ("draft-mcguinness-oauth-client-instance-assertion",
         "Extends the WG attestation framework with instance-level identity: an individual running process of a client asserts its own identity separately from the client software identity. Enables the AS to distinguish between two running instances of the same agent binary."),
        ("draft-mcguinness-oauth-ai-agent-instance",
         "Profiles client-instance-assertion specifically for AI agent platforms. Adds agent-specific metadata claims (model version, tool set, execution environment) so resource servers can make policy decisions based on the agent's runtime characteristics rather than just its software identity."),
    ]
)

# ── Slide 7: Cluster D ─────────────────────────────────────────────────────────
cluster_slide("D",
    "Actor Chain Representation",
    "No WG anchor — RFC 8693 defines the act claim but not chain validation rules or cryptographic verification",
    ORANGE,
    [
        ("draft-mcguinness-oauth-actor-profile",
         "Defines structural rules for the act claim: how to build, validate, and depth-limit delegation chains expressed as nested act objects. Establishes chain integrity invariants without requiring per-hop cryptography."),
        ("draft-mcguinness-oauth-actor-receipts",
         "Adds per-hop signed receipts to the actor chain: each actor in the delegation chain produces a verifiable receipt that proves it handled the token at that hop, enabling post-hoc audit of who touched a token and in what order."),
        ("draft-mcguinness-oauth-actor-proofs",
         "Adds hash-chained per-hop proof claims: each actor computes a proof over the previous hop's proof and its own assertion, creating a tamper-evident chain where any modification to an intermediate hop is detectable."),
        ("draft-mw-oauth-actor-chain",
         "Defines six cryptographic chain profiles organized across two axes (Declared vs. Verified binding × Full / Subset / Actor-Only chain) and directly addresses chain-splicing and chain-replay attacks. Comprehensive treatment at 97 pages."),
        ("draft-liu-oauth-chain-delegation",
         "Represents the delegation chain as a flat delegation_chain array (vs. nested act objects), with per-hop authorization constraints attached to each link. Co-authored by Aaron Parecki."),
    ]
)

# ── Slide 8: Cluster E ─────────────────────────────────────────────────────────
cluster_slide("E",
    "Token Attenuation / Sub-delegation",
    "No WG anchor — the OAuth WG recharter 'Complex Delegation' milestone is the expected adoption path",
    ORANGE,
    [
        ("draft-niyikiza-oauth-attenuating-agent-tokens",
         "Defines Attenuating Authorization Tokens (AATs): RFC 9396 RAR-based credentials that enforce monotonic narrowing — each delegation step can only remove permissions, never add. Supports offline chain verification without AS round-trips."),
        ("draft-mishra-oauth-agent-grants (DAAP)",
         "Introduces new JWT claims (agt, dev, grnt, scp, bdg) for agent delegation with depth-limiting, cascade revocation (revoking a parent token invalidates all children), and budget-bounded grants that expire on resource consumption."),
        ("draft-li-oauth-delegated-authorization",
         "Defines a hierarchical delegation token type where a client mints subordinate, narrowly-scoped access tokens for delegated parties without returning to the AS. Emphasizes client-side issuance for offline or low-latency scenarios."),
        ("draft-ni-oauth-batch-authorization-delegation",
         "Addresses multi-agent orchestration efficiency: a leader agent obtains a batch of RAR-based, actor-bound permissions in a single request and distributes subsets to sub-agents via RFC 8693 Token Exchange, reducing per-sub-agent round trips."),
    ]
)

# ── Slide 9: Cluster F ─────────────────────────────────────────────────────────
cluster_slide("F",
    "Pre-Action Authorization / Permit-Before-Commit",
    "No WG anchor — addresses the gap where standard OAuth authorizes resource access but not irreversible agent actions",
    ORANGE,
    [
        ("draft-lee-orprg-permit-receipts",
         "Abstract framework and requirements model for PermitReceipts — structured authorizations that a human or policy engine issues before an agent executes an irreversible action. Defines the data model and verifier roles without specifying a wire protocol."),
        ("draft-nelson-agent-delegation-receipts",
         "Concrete protocol: the user's private key signs Authorization Objects before runtime; results are written to an append-only log. The most iterated draft in this space at revision -10, with implementations referenced."),
        ("draft-williams-intent-token",
         "Lightweight alternative: a signed human-declared authorization envelope that travels alongside the OAuth access token. Explicitly positioned as a complement to OAuth rather than a replacement, with a minimal footprint suitable for constrained environments."),
    ]
)

# ── Slide 10: Cluster G ────────────────────────────────────────────────────────
cluster_slide("G",
    "HTTP Authorization Challenge Signals",
    "No WG anchor — draft-ietf-oauth-first-party-apps covers user step-up, a related but distinct problem",
    ORANGE,
    [
        ("draft-rosomakho-oauth-txn-challenge",
         "RS-initiated: the resource server issues a transaction-specific challenge when it determines the current authorization is insufficient for the requested operation. The client returns to the AS for human approval of that specific transaction before retrying."),
        ("draft-kahrer-oauth-client-challenge-protocol",
         "AS-initiated: the authorization server challenges the client to present additional assertions or verifiable credentials mid-flow. Addresses scenarios where the AS, not the RS, detects the insufficiency — for example when step-up policy changes after token issuance."),
        ("draft-mcguinness-oauth-insufficient-claims",
         "RS-initiated: the resource server signals precisely which credential claims are missing from the presented token, allowing the client to seek a targeted upgrade rather than a generic step-up. Complements WWW-Authenticate with a structured machine-readable claim list."),
    ]
)

# ── Slide 11: Cluster H ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s, RGBColor(0xFF, 0xF8, 0xF8))
title_bar(s, "Cluster H — AIP Name Collision",
          "Three unrelated drafts share the acronym 'AIP' — no draft can progress until re-slugged")
footer(s)

add_rect(s, 0.4, 1.25, 12.5, 0.5, fill=RED)
add_label(s,
    "The IETF secretariat must force re-slugs before these drafts can be meaningfully compared or advance.",
    0.55, 1.33, 12.2, 0.38, size=12, bold=True, color=WHITE)

aip_drafts = [
    ("draft-singla-agent-identity-protocol  (AIP)",
     "Broadest scope: a six-layer identity and authorization model spanning identity, trust establishment, permission scoping, grant issuance, consent, and reputation. Uses W3C DIDs as the identity substrate and defines three grant tiers (G1/G2/G3) for different delegation depths. 75+ pages."),
    ("draft-prakash-aip  (AIP)",
     "Technically focused: defines Intent-Bound Capability Tokens (IBCTs) using JWT/Ed25519 and Biscuit+Datalog for policy expression. Includes concrete MCP and A2A protocol bindings. Datalog-based attenuation model enables offline policy evaluation without AS contact."),
    ("draft-aip-agent-identity-protocol  (AIP)",
     "Operationally focused: introduces per-tool-call AIP Tokens, YAML-based policy definitions, and an HITL proxy pattern for human-in-the-loop approval of individual tool invocations. Significant conceptual overlap with draft-singla-agent-identity-protocol."),
]

for i, (name, desc) in enumerate(aip_drafts):
    top = 1.95 + i * 1.62
    add_rect(s, 0.4, top, 12.5, 0.38, fill=RGBColor(0xF8, 0xEC, 0xEC))
    add_label(s, name, 0.55, top + 0.05, 12.1, 0.3, size=11, bold=True, color=RED)
    add_label(s, desc, 0.55, top + 0.43, 12.1, 1.1, size=10.5, color=DGREY)

# ── Slide 12: Cluster I ────────────────────────────────────────────────────────
cluster_slide("I",
    "WIMSE Extensions for AI Agents",
    "WIMSE WG core suite (arch, identifier, mutual-tls, http-signature) provides the foundation; these drafts extend it to AI agent use cases",
    BLUE,
    [
        ("draft-ni-wimse-ai-agent-identity  (-02)",
         "Applicability statement mapping the WIMSE workload identity framework to AI agent deployments. Defines how WIMSE identifiers, credential formats, and trust establishment apply when the workload is an AI agent rather than a traditional microservice."),
        ("draft-reece-wimse-cross-org-delegation",
         "Problem statement for the cross-organization delegation gap: when an agent credential issued in one organization's WIMSE domain must be accepted by a resource in another domain, the WIMSE core provides no standard mechanism for cross-domain trust establishment."),
        ("draft-jiang-wimse-heterogeneous-credential",
         "Defines multi-format credential verification across heterogeneous WIMSE deployments where different organizations use different credential formats (SPIFFE SVIDs, JWTs, X.509). Enables interoperability without requiring format normalization upfront."),
        ("draft-schwenkschuster-wimse-trust-domain-discovery",
         "Specifies a WIMSE Trust Bundle format and a /.well-known/wimse-trust-domain discovery endpoint for obtaining cryptographic trust anchors. Companions draft-ietf-oauth-spiffe-client-auth (already WG-adopted) for the OAuth-facing side of WIMSE trust."),
        ("draft-winmagic-wimse-condition-bounded-credentials",
         "Defines credential profiles where validity is gated by attested runtime conditions (security posture, environment state) rather than a fixed expiry time. Enables agents operating in disconnected or air-gapped environments to hold credentials that remain valid only while conditions hold."),
        ("draft-munoz-wimse-authorization-evidence",
         "Defines signed authorization-evidence records that capture what a WIMSE-identified workload was authorized to do at the time of execution. Enables post-hoc audit without requiring a live AS query."),
    ]
)

# ── Save ───────────────────────────────────────────────────────────────────────
OUTPUT = "/Users/gffletch/Develop/Authorization/da_research/ietf_consolidation_analysis.pptx"
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
